import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import pandas as pd
import base64
import os
import logging
from datetime import datetime
from typing import List, Optional, Tuple
import tempfile
import json
import hashlib

# User management system
class UserManager:
    def __init__(self):
        self.users_file = "users.json"
        self.load_users()
    
    def load_users(self):
        """Load users from file or create empty dict"""
        try:
            if os.path.exists(self.users_file):
                with open(self.users_file, 'r') as f:
                    self.users = json.load(f)
            else:
                self.users = {}
        except:
            self.users = {}
    
    def save_users(self):
        """Save users to file"""
        try:
            with open(self.users_file, 'w') as f:
                json.dump(self.users, f, indent=2)
            return True
        except:
            return False
    
    def hash_password(self, password: str) -> str:
        """Hash password using SHA256"""
        return hashlib.sha256(password.encode()).hexdigest()
    
    def create_user(self, username: str, password: str, email: str = "") -> Tuple[bool, str]:
        """Create new user"""
        if username in self.users:
            return False, "Username already exists"
        
        if len(username) < 3:
            return False, "Username must be at least 3 characters"
        
        if len(password) < 6:
            return False, "Password must be at least 6 characters"
        
        self.users[username] = {
            "password": self.hash_password(password),
            "email": email,
            "created_at": datetime.now().isoformat(),
            "login_count": 0
        }
        
        if self.save_users():
            return True, "Account created successfully!"
        else:
            return False, "Error creating account. Please try again."
    
    def authenticate(self, username: str, password: str) -> Tuple[bool, str]:
        """Authenticate user"""
        if username not in self.users:
            return False, "Username not found"
        
        if self.users[username]["password"] == self.hash_password(password):
            # Update login count
            self.users[username]["login_count"] += 1
            self.users[username]["last_login"] = datetime.now().isoformat()
            self.save_users()
            return True, "Login successful!"
        else:
            return False, "Invalid password"
    
    def get_user_info(self, username: str) -> dict:
        """Get user information"""
        return self.users.get(username, {})

# Initialize user manager
user_manager = UserManager()

def enhanced_auth():
    """Enhanced authentication with sign-in and sign-up"""
    if 'authenticated' not in st.session_state:
        st.session_state.authenticated = False
        st.session_state.username = ""
    
    if not st.session_state.authenticated:
        # Custom CSS for login page
        st.markdown("""
        <style>
        .login-container {
            max-width: 450px;
            margin: 30px auto;
            padding: 2rem;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            border-radius: 20px;
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
            color: white;
        }
        
        .logo {
            text-align: center;
            margin-bottom: 30px;
        }
        
        .logo i {
            font-size: 60px;
            margin-bottom: 10px;
        }
        
        .app-title {
            font-size: 28px;
            font-weight: 700;
            margin-bottom: 8px;
            text-align: center;
        }
        
        .app-subtitle {
            font-size: 16px;
            font-weight: 300;
            text-align: center;
            opacity: 0.8;
            margin-bottom: 30px;
        }
        
        .auth-tabs {
            display: flex;
            margin-bottom: 30px;
            background: rgba(255, 255, 255, 0.1);
            border-radius: 15px;
            padding: 5px;
        }
        
        .auth-tab {
            flex: 1;
            text-align: center;
            padding: 12px;
            border-radius: 12px;
            cursor: pointer;
            transition: all 0.3s ease;
            font-weight: 600;
        }
        
        .auth-tab.active {
            background: rgba(255, 255, 255, 0.2);
            color: white;
        }
        
        .auth-tab:not(.active) {
            opacity: 0.7;
        }
        
        .stTextInput > div > div > input {
            background: rgba(255, 255, 255, 0.1);
            border: 1px solid rgba(255, 255, 255, 0.2);
            border-radius: 12px;
            color: white;
            padding: 15px;
        }
        
        .stTextInput > div > div > input::placeholder {
            color: rgba(255, 255, 255, 0.6);
        }
        
        .stButton > button {
            width: 100%;
            background: linear-gradient(45deg, #ff6b6b, #4ecdc4);
            border: none;
            border-radius: 12px;
            color: white;
            font-weight: 600;
            padding: 15px;
            text-transform: uppercase;
            letter-spacing: 1px;
            transition: all 0.3s ease;
        }
        
        .stButton > button:hover {
            transform: translateY(-2px);
            box-shadow: 0 10px 25px rgba(0, 0, 0, 0.2);
        }
        
        .features {
            text-align: center;
            margin-top: 30px;
        }
        
        .feature-icons {
            display: flex;
            justify-content: center;
            gap: 20px;
            margin-top: 15px;
        }
        
        .feature-icon {
            font-size: 20px;
            opacity: 0.7;
            transition: all 0.3s ease;
        }
        
        .feature-icon:hover {
            opacity: 1;
            transform: scale(1.2);
        }
        
        .info-box {
            background: rgba(255, 255, 255, 0.1);
            border-radius: 12px;
            padding: 15px;
            margin-top: 20px;
            text-align: center;
        }
        </style>
        
        <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css" rel="stylesheet">
        """, unsafe_allow_html=True)
        
        # Login container
        st.markdown("""
        <div class="login-container">
            <div class="logo">
                <i class="fas fa-robot"></i>
                <div class="app-title">AI PDF Chat</div>
                <div class="app-subtitle">Intelligent Document Assistant</div>
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        # Authentication tabs
        col1, col2, col3 = st.columns([1, 2, 1])
        
        with col2:
            # Tab selection
            if 'auth_tab' not in st.session_state:
                st.session_state.auth_tab = 'signin'
            
            tab_col1, tab_col2 = st.columns(2)
            
            with tab_col1:
                if st.button("🔐 Sign In", use_container_width=True, 
                           type="primary" if st.session_state.auth_tab == 'signin' else "secondary"):
                    st.session_state.auth_tab = 'signin'
                    st.rerun()
            
            with tab_col2:
                if st.button("📝 Sign Up", use_container_width=True,
                           type="primary" if st.session_state.auth_tab == 'signup' else "secondary"):
                    st.session_state.auth_tab = 'signup'
                    st.rerun()
            
            st.markdown("<br>", unsafe_allow_html=True)
            
            # Sign In Form
            if st.session_state.auth_tab == 'signin':
                st.markdown("### 🔐 Sign In to Your Account")
                
                with st.form("signin_form"):
                    username = st.text_input("👤 Username", placeholder="Enter your username")
                    password = st.text_input("🔒 Password", type="password", placeholder="Enter your password")
                    
                    signin_button = st.form_submit_button("🚀 Sign In", use_container_width=True)
                    
                    if signin_button:
                        if username and password:
                            success, message = user_manager.authenticate(username, password)
                            if success:
                                st.session_state.authenticated = True
                                st.session_state.username = username
                                st.success(message)
                                st.rerun()
                            else:
                                st.error(message)
                        else:
                            st.error("Please fill in all fields")
                
                # Info for new users
                st.markdown("""
                <div class="info-box">
                    <p><i class="fas fa-info-circle"></i> <strong>New User?</strong></p>
                    <p>Click on "Sign Up" tab above to create a new account!</p>
                </div>
                """, unsafe_allow_html=True)
            
            # Sign Up Form
            else:
                st.markdown("### 📝 Create New Account")
                
                with st.form("signup_form"):
                    new_username = st.text_input("👤 Choose Username", placeholder="Enter desired username (min 3 characters)")
                    new_email = st.text_input("📧 Email (Optional)", placeholder="Enter your email address")
                    new_password = st.text_input("🔒 Create Password", type="password", placeholder="Enter password (min 6 characters)")
                    confirm_password = st.text_input("🔒 Confirm Password", type="password", placeholder="Confirm your password")
                    
                    signup_button = st.form_submit_button("✨ Create Account", use_container_width=True)
                    
                    if signup_button:
                        if new_username and new_password and confirm_password:
                            if new_password != confirm_password:
                                st.error("Passwords don't match!")
                            else:
                                success, message = user_manager.create_user(new_username, new_password, new_email)
                                if success:
                                    st.success(message)
                                    st.info("Now you can sign in with your new account!")
                                    st.session_state.auth_tab = 'signin'
                                    st.rerun()
                                else:
                                    st.error(message)
                        else:
                            st.error("Please fill in all required fields")
                
                # Password requirements
                st.markdown("""
                <div class="info-box">
                    <p><i class="fas fa-shield-alt"></i> <strong>Account Requirements:</strong></p>
                    <p>• Username: At least 3 characters</p>
                    <p>• Password: At least 6 characters</p>
                    <p>• Email: Optional but recommended</p>
                </div>
                """, unsafe_allow_html=True)
            
            # Features showcase
            st.markdown("""
            <div class="features">
                <h4>✨ Features</h4>
                <div class="feature-icons">
                    <div class="feature-icon" title="PDF Processing">
                        <i class="fas fa-file-pdf"></i>
                    </div>
                    <div class="feature-icon" title="AI Chat">
                        <i class="fas fa-comments"></i>
                    </div>
                    <div class="feature-icon" title="Smart Search">
                        <i class="fas fa-search"></i>
                    </div>
                    <div class="feature-icon" title="Export Data">
                        <i class="fas fa-download"></i>
                    </div>
                    <div class="feature-icon" title="Secure">
                        <i class="fas fa-shield-alt"></i>
                    </div>
                </div>
            </div>
            """, unsafe_allow_html=True)
        
        st.stop()

# Call enhanced auth
enhanced_auth()

# Updated imports - using only google-generativeai directly
try:
    import google.generativeai as genai
    from sentence_transformers import SentenceTransformer
    import faiss
    import numpy as np
except ImportError as e:
    st.error(f"Missing package: {e}")
    st.info("Run: pip install google-generativeai sentence-transformers faiss-cpu numpy")
    st.stop()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PDFChatApp:
    def __init__(self):
        self.setup_page_config()
        self.initialize_session_state()
        
    def setup_page_config(self):
        """Configure Streamlit page settings"""
        st.set_page_config(
            page_title="AI-Powered PDF Chat Assistant",
            page_icon="🤖",
            layout="wide",
            initial_sidebar_state="expanded"
        )
        
    def initialize_session_state(self):
        """Initialize session state variables"""
        if 'conversation_history' not in st.session_state:
            st.session_state.conversation_history = []
        if 'vector_store' not in st.session_state:
            st.session_state.vector_store = None
        if 'processed_pdfs' not in st.session_state:
            st.session_state.processed_pdfs = []
        if 'processing_status' not in st.session_state:
            st.session_state.processing_status = "ready"
        if 'pdf_hash' not in st.session_state:
            st.session_state.pdf_hash = None
        if 'embeddings_cache' not in st.session_state:
            st.session_state.embeddings_cache = None
        if 'pdf_texts' not in st.session_state:
            st.session_state.pdf_texts = {}
            
    def get_pdf_hash(self, pdf_docs: List) -> str:
        """Generate a hash for the uploaded PDFs to enable caching"""
        hasher = hashlib.md5()
        for pdf in pdf_docs:
            pdf.seek(0)
            hasher.update(pdf.read())
            pdf.seek(0)
        return hasher.hexdigest()
    
    def is_vector_store_cached(self, pdf_docs: List) -> bool:
        """Check if vector store is already cached for these PDFs"""
        current_hash = self.get_pdf_hash(pdf_docs)
        return (st.session_state.pdf_hash == current_hash and 
                st.session_state.vector_store is not None)
    
    def validate_api_key(self, api_key: str) -> bool:
        """Validate API key format"""
        if not api_key:
            return False
        return len(api_key) > 20 and api_key.startswith('AIza')
    
    def extract_file_text(self, files: List) -> str:
        """Extract text from PDF, PPTX, and DOCX files"""
        text = ""
        successful_files = []
        failed_files = []

        for file in files:
            file_text = ""
            try:
                if file.name.lower().endswith('.pdf'):
                    pdf_reader = PdfReader(file)
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            file_text += page_text + "\n"
                elif file.name.lower().endswith('.pptx'):
                    prs = Presentation(file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                file_text += shape.text + "\n"
                elif file.name.lower().endswith('.docx'):
                    doc = Document(file)
                    for para in doc.paragraphs:
                        file_text += para.text + "\n"
                else:
                    failed_files.append(file.name)
                    continue

                if file_text.strip():
                    text += f"\n\n--- Content from {file.name} ---\n\n{file_text}"
                    successful_files.append(file.name)
                else:
                    failed_files.append(file.name)
            except Exception as e:
                failed_files.append(file.name)
                logger.error(f"Error processing {file.name}: {str(e)}")

        if failed_files:
            st.warning(f"Could not process: {', '.join(failed_files)}")
        if successful_files:
            st.success(f"Successfully processed: {', '.join(successful_files)}")
        return text
    
    def create_text_chunks(self, text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        """Split text into chunks for processing"""
        if not text.strip():
            return []
        
        # Simple text chunking
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - chunk_overlap):
            chunk = ' '.join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
        
        st.info(f"Created {len(chunks)} text chunks")
        return chunks
    
    def create_vector_store(self, text_chunks: List[str], api_key: str) -> Optional[object]:
        """Create vector store using sentence transformers and FAISS"""
        if not text_chunks:
            st.error("No text chunks to process")
            return None
            
        try:
            # Use sentence transformers for embeddings
            model = SentenceTransformer('all-MiniLM-L6-v2')
            
            with st.spinner(f"Creating vector embeddings for {len(text_chunks)} chunks..."):
                # Generate embeddings
                embeddings = model.encode(text_chunks, show_progress_bar=True)
                
                # Create FAISS index
                dimension = embeddings.shape[1]
                index = faiss.IndexFlatL2(dimension)
                index.add(embeddings.astype('float32'))
                
                # Store both index and chunks
                vector_store = {
                    'index': index,
                    'chunks': text_chunks,
                    'model': model
                }
                
            st.success(f"Vector store created successfully with {len(text_chunks)} embeddings!")
            return vector_store
            
        except Exception as e:
            st.error(f"Error creating vector store: {str(e)}")
            logger.error(f"Vector store creation failed: {str(e)}")
            return None
    
    def search_similar_chunks(self, query: str, vector_store: dict, k: int = 5) -> List[str]:
        """Search for similar chunks using the vector store"""
        try:
            # Encode the query
            query_embedding = vector_store['model'].encode([query])
            
            # Search in FAISS index
            distances, indices = vector_store['index'].search(query_embedding.astype('float32'), k)
            
            # Get relevant chunks
            relevant_chunks = []
            for idx in indices[0]:
                if idx < len(vector_store['chunks']):
                    relevant_chunks.append(vector_store['chunks'][idx])
            
            return relevant_chunks
            
        except Exception as e:
            st.error(f"Error searching chunks: {str(e)}")
            return []
    
    def generate_response(self, question: str, context: str, api_key: str) -> str:
        """Generate response using Google Gemini"""
        try:
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-1.5-flash')
            
            prompt = f"""
            Based on the following context from PDF documents, please answer the question accurately:
            
            Context:
            {context}
            
            Question: {question}
            
            Instructions:
            - Only use information from the provided context
            - If the answer is not in the context, say "The answer is not available in the provided documents"
            - Be detailed and comprehensive when possible
            - Maintain a professional tone
            
            Answer:
            """
            
            response = model.generate_content(prompt)
            return getattr(response, "text", str(response))
            
        except Exception as e:
            return f"Error generating response: {str(e)}"
    
    def process_user_question(self, user_question: str, api_key: str, files: List) -> None:
        """Process user question and generate response"""
        if not user_question.strip():
            st.warning("Please enter a question")
            return

        if not api_key or not self.validate_api_key(api_key):
            st.error("Please provide a valid API key")
            return

        if not files:
            st.error("Please upload files first")
            return

        try:
            # Check if we can use cached vector store
            if self.is_vector_store_cached(files):
                st.info("Using cached vector store for faster processing!")
                vector_store = st.session_state.vector_store
            else:
                # Process files and create new vector store
                with st.spinner("Extracting text from files..."):
                    text = self.extract_file_text(files)

                if not text.strip():
                    st.error("No text could be extracted from the uploaded files")
                    return

                with st.spinner("Processing text chunks..."):
                    text_chunks = self.create_text_chunks(text)

                if not text_chunks:
                    st.error("Could not create text chunks")
                    return

                vector_store = self.create_vector_store(text_chunks, api_key)
                if not vector_store:
                    return

                # Cache the vector store and file hash
                st.session_state.vector_store = vector_store
                st.session_state.pdf_hash = self.get_pdf_hash(files)

            # Search for relevant chunks
            with st.spinner("Searching for relevant information..."):
                relevant_chunks = self.search_similar_chunks(user_question, vector_store, k=5)

                if not relevant_chunks:
                    st.warning("No relevant content found for your question.")
                    return

                context = "\n\n".join(relevant_chunks)

            # Generate response
            with st.spinner("Generating response..."):
                response = self.generate_response(user_question, context, api_key)

            # Store conversation
            file_names = [file.name for file in files]
            timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

            st.session_state.conversation_history.append({
                "question": user_question,
                "answer": response,
                "timestamp": timestamp,
                "pdf_names": file_names,
                "model": "Google AI - Gemini 1.5 Flash",
                "user": st.session_state.username
            })

            # Display the conversation
            self.display_conversation(user_question, response, file_names)

        except Exception as e:
            st.error(f"Error processing question: {str(e)}")
            logger.error(f"Question processing failed: {str(e)}")
    
    def display_conversation(self, question: str, answer: str, pdf_names: List[str]):
        """Display conversation in a chat-like interface"""
        st.markdown("### 💬 Conversation")
        
        # Current conversation
        st.markdown(
            f"""
            <div style="background-color: #f0f2f6; padding: 15px; border-radius: 10px; margin: 10px 0;">
                <div style="display: flex; align-items: center; margin-bottom: 10px;">
                    <div style="background-color: #4CAF50; color: white; padding: 5px 10px; border-radius: 15px; font-size: 12px;">
                        {st.session_state.username.upper()}
                    </div>
                </div>
                <div style="margin-left: 15px; font-size: 16px; color: #222;">
                    {question}
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )
        
        st.markdown(
            f"""
            <div style="background-color: #e3f2fd; padding: 15px; border-radius: 10px; margin: 10px 0;">
                <div style="display: flex; align-items: center; margin-bottom: 10px;">
                    <div style="background-color: #2196F3; color: white; padding: 5px 10px; border-radius: 15px; font-size: 12px;">
                        AI ASSISTANT
                    </div>
                </div>
                <div style="margin-left: 15px; font-size: 16px; line-height: 1.6; color: #222;">
                    {answer}
                </div>
                <div style="margin-top: 10px; font-size: 12px; color: #666;">
                    📄 Sources: {', '.join(pdf_names)}
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )
        
        # Display conversation history
        if len(st.session_state.conversation_history) > 1:
            st.markdown("### 📝 Previous Conversations")
            for i, conv in enumerate(reversed(st.session_state.conversation_history[:-1])):
                with st.expander(f"Conversation {len(st.session_state.conversation_history) - i - 1}: {conv['question'][:50]}..."):
                    st.write(f"**Question:** {conv['question']}")
                    st.write(f"**Answer:** {conv['answer']}")
                    st.write(f"**Time:** {conv['timestamp']}")
                    st.write(f"**Sources:** {', '.join(conv['pdf_names'])}")
                    if 'user' in conv:
                        st.write(f"**User:** {conv['user']}")
    
    def export_conversation_history(self):
        """Export conversation history to various formats"""
        if not st.session_state.conversation_history:
            return
            
        export_data = []
        for conv in st.session_state.conversation_history:
            export_data.append({
                "Question": conv["question"],
                "Answer": conv["answer"],
                "Timestamp": conv["timestamp"],
                "PDF Sources": ", ".join(conv["pdf_names"]),
                "Model": conv["model"],
                "User": conv.get("user", "Unknown")
            })
        
        df = pd.DataFrame(export_data)
        
        # CSV Export
        csv = df.to_csv(index=False)
        st.sidebar.download_button(
            label="📊 Download CSV",
            data=csv,
            file_name=f"pdf_chat_history_{st.session_state.username}.csv",
            mime="text/csv"
        )
        
        # JSON Export
        json_data = json.dumps(export_data, indent=2)
        st.sidebar.download_button(
            label="📋 Download JSON",
            data=json_data,
            file_name=f"pdf_chat_history_{st.session_state.username}.json",
            mime="application/json"
        )
    
    def render_sidebar(self):
        """Render the sidebar with controls"""
        st.sidebar.title("🤖 AI PDF Chat Assistant")
        st.sidebar.markdown("---")
        
        # User info and logout
        user_info = user_manager.get_user_info(st.session_state.username)
        st.sidebar.success(f"✅ Welcome **{st.session_state.username}**!")
        
        if user_info:
            login_count = user_info.get('login_count', 0)
            st.sidebar.info(f"🔢 Total logins: {login_count}")
            
            if 'created_at' in user_info:
                created_date = datetime.fromisoformat(user_info['created_at']).strftime('%Y-%m-%d')
                st.sidebar.info(f"📅 Member since: {created_date}")
        
        if st.sidebar.button("🚪 Logout"):
            st.session_state.authenticated = False
            st.session_state.username = ""
            # Clear session data on logout
            for key in list(st.session_state.keys()):
                if key not in ['authenticated', 'username']:
                    del st.session_state[key]
            st.rerun()
        
        st.sidebar.markdown("---")
        
        # Profile links
        st.sidebar.markdown("### 👨‍💻 Developer")
        linkedin_link = "https://www.linkedin.com/in/sairakshith-talluru-a69272265/"
        kaggle_link = "https://www.kaggle.com/sairakshith2004"
        github_link = "https://github.com/sairakshith2004"
        
        st.sidebar.markdown(f"[LinkedIn]({linkedin_link})")
        st.sidebar.markdown(f"[Kaggle]({kaggle_link})")
        st.sidebar.markdown(f"[GitHub]({github_link})")
        
        st.sidebar.markdown("---")
        
        # API Key Configuration
        st.sidebar.markdown("### 🔑 API Configuration")
        api_key = st.sidebar.text_input(
            "Enter your Google AI API Key:",
            type="password",
            help="Get your API key from Google AI Studio"
        )
        
        if api_key and not self.validate_api_key(api_key):
            st.sidebar.error("⚠️ Invalid API key format")
        elif api_key:
            st.sidebar.success("✅ API key configured")
        
        st.sidebar.markdown("---")
        
        # Statistics
        if st.session_state.conversation_history:
            st.sidebar.markdown("### 📊 Session Statistics")
            st.sidebar.metric("Total Questions", len(st.session_state.conversation_history))
            
            if st.session_state.processed_pdfs:
                st.sidebar.metric("PDFs Processed", len(st.session_state.processed_pdfs))
        
        # Export options
        if st.session_state.conversation_history:
            st.sidebar.markdown("### 📥 Export Options")
            self.export_conversation_history()
        
        # Clear conversation
        if st.sidebar.button("🗑️ Clear Conversation"):
            st.session_state.conversation_history = []
            st.sidebar.success("Conversation cleared!")
            st.rerun()
        
        st.sidebar.markdown("---")
        st.sidebar.markdown("### ℹ️ Instructions")
        st.sidebar.info(
            """
            1. Enter your Google AI API key
            2. Upload PDF files
            3. Ask questions about the content
            4. Export your conversation history
            """
        )
        
        return api_key
    
    def render_main_content(self, api_key: str):
        """Render the main content area"""
        # Header
        st.markdown(
            """
            <div style="text-align: center; padding: 20px; background: linear-gradient(90deg, #667eea 0%, #764ba2 100%); border-radius: 10px; margin-bottom: 20px;">
                <h1 style="color: white; margin: 0;">🤖 AI-Powered PDF Chat Assistant</h1>
                <p style="color: white; margin: 5px 0 0 0; opacity: 0.9;">Upload your PDFs and chat with your documents using AI</p>
            </div>
            """,
            unsafe_allow_html=True
        )
        
        # File upload section
        st.markdown("### 📄 Upload Documents (PDF, PPTX, DOCX)")
        files = st.file_uploader(
            "Choose files",
            accept_multiple_files=True,
            type=["pdf", "pptx", "docx"],
            help="You can upload multiple files (PDF, PPTX, DOCX) at once"
        )

        if files is not None and len(files) > 0:
            # Display uploaded files info
            col1, col2, col3 = st.columns(3)

            with col1:
                st.metric("Files Uploaded", len(files))

            with col2:
                total_size = sum([file.size for file in files]) / (1024 * 1024)  # MB
                st.metric("Total Size", f"{total_size:.2f} MB")

            with col3:
                if st.session_state.vector_store and self.is_vector_store_cached(files):
                    st.metric("Status", "✅ Cached")
                else:
                    st.metric("Status", "🔄 Ready to Process")

            # Show file details
            with st.expander("📋 File Details"):
                for idx, file in enumerate(files, 1):
                    size_mb = file.size / (1024 * 1024)
                    st.write(f"**{idx}.** {file.name} ({size_mb:.2f} MB)")

            # Automatically explain the uploaded documents using AI
            if api_key and self.validate_api_key(api_key):
                self.explain_documents(files, api_key)

        st.markdown("---")
        
        # Chat interface
        st.markdown("### 💬 Chat with Your Documents")
        
        # Question input
        user_question = st.text_input(
            "Ask a question about your documents:",
            placeholder="e.g., What is the main topic discussed in the document?",
            help="Type your question and press Enter"
        )
        
        # Submit button with enhanced styling
        col1, col2, col3 = st.columns([2, 1, 2])
        with col2:
            submit_button = st.button(
                "🚀 Ask Question",
                use_container_width=True,
                type="primary"
            )
        
        # Process question when submitted
        if submit_button or user_question:
            if user_question:
                self.process_user_question(user_question, api_key, files)
        
        # Display existing conversation history if available
        if st.session_state.conversation_history and not (submit_button or user_question):
            st.markdown("### 📖 Recent Conversations")
            
            # Show last few conversations
            recent_conversations = st.session_state.conversation_history[-3:]
            
            for i, conv in enumerate(reversed(recent_conversations)):
                with st.expander(f"💭 {conv['question'][:60]}..." if len(conv['question']) > 60 else f"💭 {conv['question']}"):
                    st.markdown(f"**🙋 Question:** {conv['question']}")
                    st.markdown(f"**🤖 Answer:** {conv['answer']}")
                    st.caption(f"⏰ {conv['timestamp']} | 📄 {', '.join(conv['pdf_names'])} | 👤 {conv.get('user', 'Unknown')}")
        
        # Footer
        st.markdown("---")
        st.markdown(
            """
            <div style="text-align: center; padding: 10px; color: #666;">
                <p>Built with ❤️ using Streamlit, Google AI, and advanced NLP techniques</p>
                <p><small>Secure • Private • Intelligent</small></p>
            </div>
            """,
            unsafe_allow_html=True
        )
    
    def run(self):
        """Main application runner"""
        try:
            # Render sidebar and get API key
            api_key = self.render_sidebar()
            
            # Render main content
            self.render_main_content(api_key)
            
        except Exception as e:
            st.error(f"Application error: {str(e)}")
            logger.error(f"Application error: {str(e)}")

    def explain_documents(self, files: List, api_key: str):
        """Generate an AI-powered summary/explanation of the uploaded documents"""
        if not files or not api_key or not self.validate_api_key(api_key):
            return

        with st.spinner("AI is analyzing your documents..."):
            text = self.extract_file_text(files)
            if not text.strip():
                st.warning("No text could be extracted for explanation.")
                return

            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-1.5-flash')

            prompt = f"""
            Please read the following content extracted from uploaded documents (PDF, PPTX, DOCX) and provide a clear, concise, and detailed explanation about the main topics, purpose, and key points covered. If possible, summarize the content and highlight important sections.

            Content:
            {text}

            Instructions:
            - Give a summary and explanation suitable for someone unfamiliar with the document.
            - Highlight main topics, purpose, and any important details.
            - Be clear, professional, and informative.
            """

            try:
                response = model.generate_content(prompt)
                summary = getattr(response, "text", str(response))
                st.markdown(
                    f"""
                    <div style="background-color: #fffbe6; padding: 18px; border-radius: 12px; margin: 18px 0; border: 1px solid #ffe58f;">
                        <div style="font-weight: bold; color: #ad8b00; margin-bottom: 8px;">🧠 AI Document Summary</div>
                        <div style="color: #222; font-size: 16px; line-height: 1.7;">{summary}</div>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            except Exception as e:
                st.error(f"AI explanation failed: {str(e)}")

# Custom CSS for enhanced styling
def apply_custom_css():
    st.markdown("""
    <style>
    /* Main app styling */
    .main .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
    }
    
    /* Sidebar styling */
    .css-1d391kg {
        background-color: #f8f9fa;
    }
    
    /* Metric cards */
    .css-1r6slb0 {
        background-color: #ffffff;
        border: 1px solid #e0e0e0;
        border-radius: 10px;
        padding: 1rem;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    
    /* File uploader */
    .css-1cpxqw2 {
        border: 2px dashed #cccccc;
        border-radius: 10px;
        padding: 2rem;
        text-align: center;
        background-color: #fafafa;
    }
    
    /* Buttons */
    .stButton > button {
        border-radius: 20px;
        border: none;
        font-weight: 600;
        transition: all 0.3s ease;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0,0,0,0.2);
    }
    
    /* Success/Error messages */
    .stSuccess {
        border-radius: 10px;
    }
    
    .stError {
        border-radius: 10px;
    }
    
    .stWarning {
        border-radius: 10px;
    }
    
    .stInfo {
        border-radius: 10px;
    }
    
    /* Text input */
    .stTextInput > div > div > input {
        border-radius: 20px;
        border: 2px solid #e0e0e0;
        padding: 10px 15px;
    }
    
    .stTextInput > div > div > input:focus {
        border-color: #667eea;
        box-shadow: 0 0 0 0.2rem rgba(102, 126, 234, 0.25);
    }
    
    /* Expanders */
    .streamlit-expanderHeader {
        border-radius: 10px;
        background-color: #f8f9fa;
    }
    
    /* Progress bars */
    .stProgress .st-bo {
        border-radius: 10px;
    }
    
    /* Conversation styling */
    .conversation-container {
        max-height: 400px;
        overflow-y: auto;
        padding: 10px;
        border: 1px solid #e0e0e0;
        border-radius: 10px;
        background-color: #fafafa;
    }
    
    /* Hide Streamlit footer */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    
    /* Custom animations */
    @keyframes fadeIn {
        from {opacity: 0; transform: translateY(20px);}
        to {opacity: 1; transform: translateY(0);}
    }
    
    .fade-in {
        animation: fadeIn 0.5s ease-in;
    }
    </style>
    """, unsafe_allow_html=True)

# Main execution
def main():
    """Main function to run the application"""
    try:
        # Apply custom CSS
        apply_custom_css()
        
        # Initialize and run the app
        app = PDFChatApp()
        app.run()
        
    except Exception as e:
        st.error(f"Failed to start application: {str(e)}")
        logger.error(f"Application startup failed: {str(e)}")

if __name__ == "__main__":
    main()